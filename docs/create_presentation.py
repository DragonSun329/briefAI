"""
briefAI Technical Presentation Generator
Creates a comprehensive PPTX presentation with quantitative metrics and technical details.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme
DARK_BLUE = RGBColor(0x1a, 0x36, 0x5d)
LIGHT_BLUE = RGBColor(0x2e, 0x86, 0xab)
ACCENT_GREEN = RGBColor(0x28, 0xa7, 0x45)
ACCENT_ORANGE = RGBColor(0xf5, 0x7c, 0x00)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf0, 0xf0, 0xf0)


def add_title_slide(prs, title, subtitle):
    """Add a title slide with dark background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullet_points, notes=None):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(point, tuple):
            # Sub-bullet
            p.text = f"    • {point[1]}"
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.text = f"• {point}"
            p.font.size = Pt(20)

        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(12)

    return slide


def add_metrics_slide(prs, title, metrics):
    """Add a slide with large metric callouts."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Metric boxes - 2x2 or 3x2 grid
    cols = 3 if len(metrics) > 4 else 2
    rows = (len(metrics) + cols - 1) // cols

    box_width = Inches(2.8)
    box_height = Inches(1.8)
    start_x = Inches(0.6)
    start_y = Inches(1.6)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)

    for i, (value, label) in enumerate(metrics):
        row = i // cols
        col = i % cols

        x = start_x + col * (box_width + gap_x)
        y = start_y + row * (box_height + gap_y)

        # Box background
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()

        # Value
        val_box = slide.shapes.add_textbox(x, y + Inches(0.2), box_width, Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(x, y + Inches(1.1), box_width, Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a data table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.size = Pt(14)

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_BLUE

    return slide


def create_briefai_presentation():
    """Create the complete briefAI technical presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title ==========
    add_title_slide(
        prs,
        "briefAI",
        "AI-Powered Trend Intelligence Platform\nTechnical Overview & Quantitative Analysis"
    )

    # ========== SLIDE 2: Executive Summary ==========
    add_metrics_slide(prs, "Platform Overview", [
        ("236", "Python Files"),
        ("27,600+", "Lines of Code"),
        ("70+", "Data Sources"),
        ("7", "Signal Types"),
        ("31", "AI Trend Buckets"),
        ("5", "Alert Categories"),
    ])

    # ========== SLIDE 3: Section - Architecture ==========
    add_section_slide(prs, "System Architecture")

    # ========== SLIDE 4: Core Modules ==========
    add_content_slide(prs, "Core Module Structure", [
        "Data Collection Layer (scrapers/, 12 specialized scrapers)",
        ("sub", "GitHub Enhanced Scraper: 50+ repos/bucket, star velocity tracking"),
        ("sub", "HuggingFace Scraper: Model downloads, trending models"),
        ("sub", "Financial Signals: Yahoo Finance, SEC EDGAR integration"),
        "Processing Layer (utils/, 45+ utility modules)",
        ("sub", "NLP Pipeline: spaCy + Sentence Transformers for entity extraction"),
        ("sub", "Trend Aggregator: Multi-source signal fusion"),
        "Analytics Layer (modules/, 18 dashboard modules)",
        ("sub", "Streamlit-based interactive dashboards"),
        ("sub", "Real-time visualization with Plotly"),
    ])

    # ========== SLIDE 5: Data Pipeline ==========
    add_table_slide(prs, "Data Pipeline Specifications",
        ["Component", "Technology", "Update Frequency", "Volume"],
        [
            ["News Aggregation", "NewsAPI + RSS", "Every 4 hours", "500+ articles/day"],
            ["GitHub Trends", "GitHub API + Scraping", "Daily", "1,500+ repos tracked"],
            ["Financial Data", "Yahoo Finance API", "Real-time", "100+ tickers"],
            ["SEC Filings", "EDGAR API", "Daily", "10-K, 10-Q, 8-K, S-1"],
            ["VC/Funding", "Crunchbase + News", "Daily", "50+ deals/week"],
            ["HuggingFace", "HF API", "Daily", "500+ models tracked"],
        ]
    )

    # ========== SLIDE 6: Section - Signal Framework ==========
    add_section_slide(prs, "7-Signal Intelligence Framework")

    # ========== SLIDE 7: Signal Types ==========
    add_table_slide(prs, "Signal Type Specifications",
        ["Signal", "Full Name", "Data Sources", "Score Range"],
        [
            ["TMS", "Technical Momentum Score", "GitHub, HuggingFace, PapersWithCode", "0-100"],
            ["CCS", "Capital Conviction Score", "Crunchbase, PitchBook, News", "0-100"],
            ["EIS", "Enterprise Institutional Signal", "SEC 10-K, 10-Q, S-1, 8-K", "0-100"],
            ["NAS", "Narrative Attention Score", "TechCrunch, Bloomberg, VentureBeat", "0-100"],
            ["PMS", "Public Market Signal", "Yahoo Finance, Stock mentions", "0-100"],
            ["CSS", "Crypto Sentiment Signal", "Kraken, CoinGecko", "0-100"],
            ["MRS", "Macro Regime Signal", "Economic indicators", "0-100"],
        ]
    )

    # ========== SLIDE 8: Signal Computation ==========
    add_content_slide(prs, "Signal Computation Methodology", [
        "Percentile Ranking: All signals normalized to 0-100 scale",
        ("sub", "Uses rolling 4-week baseline for comparison"),
        ("sub", "Z-score normalization with 1.96 confidence interval"),
        "Velocity Tracking: Week-over-week delta computation",
        ("sub", "Velocity spike threshold: 3.0x baseline"),
        ("sub", "EWMA smoothing with α=0.3"),
        "Composite Scoring: Weighted signal aggregation",
        ("sub", "Source credibility tiers: 1.0, 0.8, 0.6, 0.4"),
        ("sub", "Minimum 2 sources required for confidence"),
    ])

    # ========== SLIDE 9: Section - AI Trend Buckets ==========
    add_section_slide(prs, "31 AI Trend Buckets")

    # ========== SLIDE 10: Bucket Categories ==========
    add_table_slide(prs, "Trend Bucket Taxonomy",
        ["Category", "Buckets", "Count"],
        [
            ["Foundation Models", "LLM Architectures, Multimodal AI, Reasoning Models", "3"],
            ["AI Infrastructure", "AI Chips, Vector DBs, Model Serving, MLOps", "4"],
            ["Agent Systems", "Agentic AI, Browser Agents, Code Agents", "3"],
            ["Enterprise AI", "AI Security, Compliance AI, Document AI", "3"],
            ["Industry Vertical", "Healthcare AI, Legal AI, Finance AI, Climate AI", "4"],
            ["Creative AI", "Text-to-Video, Music Gen, 3D Generation", "3"],
            ["Emerging Tech", "Robotics, Quantum ML, Neuro-symbolic AI", "3"],
            ["Developer Tools", "AI IDEs, Prompt Engineering, Fine-tuning", "3"],
            ["Data & Analytics", "Synthetic Data, Data Labeling, AutoML", "3"],
            ["Edge & Mobile", "On-device AI, TinyML", "2"],
        ]
    )

    # ========== SLIDE 11: Bucket Lifecycle ==========
    add_content_slide(prs, "Bucket Lifecycle Classification", [
        "Emerging (Score: 90): New technology, high TMS, low CCS/EIS",
        ("sub", "Detection: TMS > 80, CCS < 20, minimal enterprise mentions"),
        "Validating (Score: 80): Growing attention, initial funding",
        ("sub", "Detection: Rising NAS, first CCS signals appearing"),
        "Establishing (Score: 60): Enterprise pilots, sustained funding",
        ("sub", "Detection: EIS > 30, CCS > 50, stable NAS"),
        "Mainstream (Score: 40): Broad adoption, commoditizing",
        ("sub", "Detection: EIS > 60, stable all signals, low velocity"),
    ])

    # ========== SLIDE 12: Section - Business Analyst Frameworks ==========
    add_section_slide(prs, "Business Analyst Frameworks")

    # ========== SLIDE 13: Gartner Hype Cycle ==========
    add_table_slide(prs, "Gartner Hype Cycle Integration",
        ["Phase", "Signal Pattern", "Threshold"],
        [
            ["Innovation Trigger", "High TMS, Low CCS/EIS", "TMS>80, CCS<20, EIS<10"],
            ["Peak of Expectations", "Max NAS, High CCS, Flat TMS", "NAS>80, CCS>70, TMS Δ≤0"],
            ["Trough of Disillusionment", "Declining NAS/CCS", "NAS Δ<-20, CCS Δ<-10"],
            ["Slope of Enlightenment", "Rising EIS, Stable TMS", "EIS Δ>10, TMS Δ>0"],
            ["Plateau of Productivity", "All signals stable, High EIS", "Variance<10, EIS>60"],
        ]
    )

    # ========== SLIDE 14: 5T Investment Framework ==========
    add_table_slide(prs, "5T Investment Thesis Scoring",
        ["Dimension", "Weight", "Data Sources", "Key Indicators"],
        [
            ["Team", "20%", "Crunchbase, News", "Founder mentions, key hires, CB rank"],
            ["Technology", "20%", "GitHub, HuggingFace", "Star velocity, downloads, repo count"],
            ["Market", "20%", "News, Keywords", "TAM mentions, market size articles"],
            ["Timing", "20%", "Lifecycle, Velocity", "Hype cycle phase, signal deltas"],
            ["Traction", "20%", "CCS, EIS, PMS", "Funding rounds, SEC filings, revenue"],
        ]
    )

    # ========== SLIDE 15: Confidence Framework ==========
    add_content_slide(prs, "Confidence & Credibility System", [
        "Source Credibility Tiers (4-tier weighting)",
        ("sub", "Tier 1 (1.0): SEC filings, Sequoia, a16z, WSJ, Bloomberg"),
        ("sub", "Tier 2 (0.8): Crunchbase, PitchBook, Forbes, Fortune"),
        ("sub", "Tier 3 (0.6): GitHub Trending, HuggingFace, TechCrunch"),
        ("sub", "Tier 4 (0.4): Twitter/X, Reddit, Medium, Crypto sources"),
        "Confidence Intervals (95% CI)",
        ("sub", "Coverage variance: Based on entity count vs expected"),
        ("sub", "Source variance: Standard deviation across sources"),
        ("sub", "Temporal variance: Week-over-week volatility"),
        "Data Coverage Badges",
        ("sub", "FULL: 6/6 signals | GOOD: 4-5/6 | PARTIAL: 2-3/6 | LOW: 0-1/6"),
    ])

    # ========== SLIDE 16: Section - Alert System ==========
    add_section_slide(prs, "Intelligent Alert System")

    # ========== SLIDE 17: Alert Types ==========
    add_table_slide(prs, "Alert Type Specifications",
        ["Alert Type", "Trigger Condition", "Signal Pattern"],
        [
            ["Alpha Zone", "TMS>80, CCS<50, EIS<30", "Early technical momentum, undervalued"],
            ["Hype Zone", "NAS>80, TMS<60", "High attention, weak fundamentals"],
            ["Enterprise Pull", "EIS>70, TMS>60", "Strong institutional adoption"],
            ["Disruption Pressure", "TMS>90, NAS>70, CCS rising", "Rapid technical + narrative growth"],
            ["Rotation Signal", "Velocity reversal detected", "Trend direction change"],
        ]
    )

    # ========== SLIDE 18: Section - Performance ==========
    add_section_slide(prs, "System Performance Metrics")

    # ========== SLIDE 19: Performance Stats ==========
    add_metrics_slide(prs, "Processing Performance", [
        ("<4 hrs", "Data Refresh Cycle"),
        ("500+", "Articles/Day Processed"),
        ("1,500+", "GitHub Repos Tracked"),
        ("31", "Buckets Analyzed"),
        ("95%", "Uptime Target"),
        ("< 2s", "Dashboard Load Time"),
    ])

    # ========== SLIDE 20: Data Quality ==========
    add_content_slide(prs, "Data Quality Assurance", [
        "Entity Extraction: spaCy NER + custom AI taxonomy",
        ("sub", "Precision optimized for technology terms"),
        ("sub", "Company/product disambiguation pipeline"),
        "Deduplication: Fuzzy matching with 85% threshold",
        ("sub", "Cross-source entity resolution"),
        ("sub", "Temporal aggregation (weekly buckets)"),
        "Validation: Multi-source confirmation",
        ("sub", "Minimum 2 sources for trend confirmation"),
        ("sub", "3-week minimum history for baseline"),
    ])

    # ========== SLIDE 21: Section - Tech Stack ==========
    add_section_slide(prs, "Technology Stack")

    # ========== SLIDE 22: Tech Stack Details ==========
    add_table_slide(prs, "Technology Stack",
        ["Layer", "Technology", "Purpose"],
        [
            ["Language", "Python 3.11+", "Core implementation"],
            ["Web Framework", "Streamlit", "Interactive dashboards"],
            ["Visualization", "Plotly, Matplotlib", "Charts and graphs"],
            ["NLP", "spaCy, Sentence Transformers", "Entity extraction, embeddings"],
            ["Data", "Pandas, NumPy", "Data processing"],
            ["APIs", "Requests, aiohttp", "Data collection"],
            ["Storage", "JSON, SQLite", "Data persistence"],
            ["Config", "Pydantic, JSON", "Configuration management"],
        ]
    )

    # ========== SLIDE 23: File Structure ==========
    add_content_slide(prs, "Codebase Structure (236 files, 27,600+ LOC)", [
        "scrapers/ - 12 data collection modules",
        ("sub", "github_enhanced_scraper.py, huggingface_scraper.py, etc."),
        "modules/ - 18 dashboard and analysis modules",
        ("sub", "bucket_dashboard.py, trend_aggregator.py, etc."),
        "utils/ - 45+ utility modules",
        ("sub", "bucket_models.py, bucket_scorers.py, investment_thesis.py"),
        "config/ - Configuration files",
        ("sub", "trend_detection.json, five_t_scoring.json, source_credibility.json"),
        "data/ - Generated data and caches",
        ("sub", "bucket_profiles.json, trend_heatmaps/, alerts/"),
    ])

    # ========== SLIDE 24: Section - Roadmap ==========
    add_section_slide(prs, "Future Enhancements")

    # ========== SLIDE 25: Roadmap ==========
    add_content_slide(prs, "Development Roadmap", [
        "Phase 1: Enhanced Signal Coverage",
        ("sub", "Add patent filing signals (USPTO integration)"),
        ("sub", "Academic paper velocity (arXiv, Semantic Scholar)"),
        "Phase 2: Predictive Capabilities",
        ("sub", "ML-based trend prediction (3-6 month forecast)"),
        ("sub", "Similar trend pattern matching"),
        "Phase 3: Integration & APIs",
        ("sub", "REST API for programmatic access"),
        ("sub", "Slack/Teams bot integration"),
        "Phase 4: Advanced Analytics",
        ("sub", "Competitive landscape mapping"),
        ("sub", "Investment portfolio optimization"),
    ])

    # ========== SLIDE 26: Closing ==========
    add_title_slide(
        prs,
        "briefAI",
        "Transforming AI Trend Analysis with\nData-Driven Intelligence"
    )

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "..", "briefAI_Technical_Presentation.pptx")
    output_path = os.path.abspath(output_path)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_briefai_presentation()
